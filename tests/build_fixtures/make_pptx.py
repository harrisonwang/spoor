#!/usr/bin/env python3
"""PPTX fixtures."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent.parent / "fixtures" / "pptx"
OUT.mkdir(parents=True, exist_ok=True)


# ---------- 01: basic title slides ----------
def build_01_basic():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Title Slide"
    s.placeholders[1].text = "A subtitle"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Second slide"
    tf = s.placeholders[1].text_frame
    tf.text = "First bullet"
    p = tf.add_paragraph(); p.text = "Second bullet"
    p = tf.add_paragraph(); p.text = "Nested item"; p.level = 1

    prs.save(OUT / "01_basic.pptx")


# ---------- 02: with table ----------
def build_02_with_table():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Has a Table"
    rows, cols = 3, 2
    tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Header A"
    tbl.cell(0, 1).text = "Header B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    tbl.cell(2, 0).text = "3"
    tbl.cell(2, 1).text = "4"
    prs.save(OUT / "02_with_table.pptx")


# ---------- 03: speaker notes ----------
def build_03_with_notes():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide with notes"
    s.notes_slide.notes_text_frame.text = "These are speaker notes that explain the slide."
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide without notes"
    prs.save(OUT / "03_with_notes.pptx")


# ---------- 04: empty deck (one blank slide) ----------
def build_04_empty():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(OUT / "04_empty.pptx")


# ---------- 05: many slides (ordering test - ensure slide11 > slide2) ----------
def build_05_ordering():
    prs = Presentation()
    for i in range(1, 13):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"Slide number {i}"
    prs.save(OUT / "05_ordering.pptx")


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("build_") and callable(fn):
            print(f"Building {name}...")
            fn()
