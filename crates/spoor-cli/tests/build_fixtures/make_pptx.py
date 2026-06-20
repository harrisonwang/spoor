#!/usr/bin/env python3
"""PPTX fixtures."""
import base64
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent.parent / "fixtures" / "pptx"
OUT.mkdir(parents=True, exist_ok=True)


# ---------- 01: basic title slides ----------
def build_01_basic():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Title Slide"
    s.placeholders[1].text = "A subtitle"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Second slide"
    tf = s.placeholders[1].text_frame
    tf.text = "First bullet"
    p = tf.add_paragraph(); p.text = "Second bullet"
    p = tf.add_paragraph(); p.text = "Nested item"; p.level = 1

    prs.save(OUT / "01_basic.pptx")


# ---------- 02: with table ----------
def build_02_with_table():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Has a Table"
    rows, cols = 3, 2
    tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Header A"
    tbl.cell(0, 1).text = "Header B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    tbl.cell(2, 0).text = "3"
    tbl.cell(2, 1).text = "4"
    prs.save(OUT / "02_with_table.pptx")


# ---------- 03: speaker notes ----------
def build_03_with_notes():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide with notes"
    s.notes_slide.notes_text_frame.text = "These are speaker notes that explain the slide."
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide without notes"
    prs.save(OUT / "03_with_notes.pptx")


# ---------- 04: empty deck (one blank slide) ----------
def build_04_empty():
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(OUT / "04_empty.pptx")


# ---------- 05: many slides (ordering test - ensure slide11 > slide2) ----------
def build_05_ordering():
    prs = Presentation()
    for i in range(1, 13):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"Slide number {i}"
    prs.save(OUT / "05_ordering.pptx")


# ---------- 06: merged table must emit an integrity warning ----------
def build_06_merged_table():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Merged table"
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Merged header"
    tbl.cell(0, 0).merge(tbl.cell(0, 1))
    tbl.cell(1, 0).text = "A"
    tbl.cell(1, 1).text = "B"
    prs.save(OUT / "06_merged_table.pptx")


# ---------- 07: omitted picture must emit an integrity warning ----------
def build_07_embedded_visual():
    png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
        "+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Picture slide"
    s.shapes.add_picture(BytesIO(png), Inches(1), Inches(2), Inches(2), Inches(2))
    prs.save(OUT / "07_embedded_visual.pptx")


# ---------- 08: multi-slide image placeholders ----------
# Drives the `spoor://pptx/part/ppt/media/*` emission path: covers per-slide
# image numbering, multiple images on a single slide, and an image-free slide
# that should produce no handles.
def build_08_image_placeholders():
    # Two distinct 1x1 PNGs so python-pptx writes two separate media parts
    # instead of deduplicating by content hash.
    red_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+hHgAHggJ/PchI7wAAAABJRU5ErkJggg=="
    )
    blue_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNgYPj/HwADAgH/eL5gtgAAAABJRU5ErkJggg=="
    )
    prs = Presentation()
    # Slide 1: one image.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide one"
    s.shapes.add_picture(BytesIO(red_png), Inches(1), Inches(2), Inches(2), Inches(2))
    # Slide 2: two images, distinct bytes so they land in distinct media parts.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide two"
    s.shapes.add_picture(BytesIO(red_png), Inches(1), Inches(2), Inches(2), Inches(2))
    s.shapes.add_picture(BytesIO(blue_png), Inches(4), Inches(2), Inches(2), Inches(2))
    # Slide 3: title only — must produce no handle.
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Slide three (no images)"
    prs.save(OUT / "08_image_placeholders.pptx")


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("build_") and callable(fn):
            print(f"Building {name}...")
            fn()
